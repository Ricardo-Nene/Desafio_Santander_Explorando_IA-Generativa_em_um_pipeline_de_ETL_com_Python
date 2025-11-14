
import os
import sqlite3
from datetime import datetime

import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt

#CONFIG. BASICA

RAW_PATH = "data/raw/clientes.csv"
PROCESSED_DIR = "data/processed"
DB_PATH = "data/banco_do_ricardo_seguros.db"

os.makedirs(PROCESSED_DIR, exist_ok=True)

# EXTRACT (E do ETL)

def extract_clientes(path: str) -> pd.DataFrame:
    df = pd.read_csv(path)
    df["idade"] = pd.to_numeric(df["idade"], errors="coerce")
    df["renda_mensal"] = pd.to_numeric(df["renda_mensal"], errors="coerce")
    df["saldo_medio"] = pd.to_numeric(df["saldo_medio"], errors="coerce")
    return df

# REGRAS 

def classificar_faixa_renda(renda):
    if renda < 3000:
        return "baixa"
    elif renda < 8000:
        return "media"
    else:
        return "alta"


def recomendar_tipo_seguro(linha):
    idade = linha["idade"]
    renda = linha["renda_mensal"]
    saldo = linha["saldo_medio"]
    possui_seguro = linha["possui_seguro"]

    # Já tem seguro -> sugestão de upgrade
    if possui_seguro == "S":
        return "upgrade_plano_atual"

    # Regras simples de exemplo:
    if idade >= 40 and renda >= 5000:
        return "seguro_vida"
    if saldo >= 15000:
        return "seguro_previdencia"
    if 4000 <= renda < 8000:
        return "seguro_residencial"

    # fallback
    return "seguro_cartao"

# Etapa TRANSFORM (T do ETL)

def transform(df: pd.DataFrame):
    """
    Aplica regras de negócio:
    - faixa de renda
    - tipo de seguro recomendado
    - prioridade da oferta
    - resumo gerencial
    """
    df = df.copy()

    # Faixa de renda
    df["faixa_renda"] = df["renda_mensal"].apply(classificar_faixa_renda)

    # Tipo de seguro recomendado
    df["seguro_recomendado"] = df.apply(recomendar_tipo_seguro, axis=1)

    # Prioridade comercial
    def prioridade(row):
        if row["seguro_recomendado"] in ["seguro_vida", "seguro_previdencia"]:
            return "alta"
        if row["seguro_recomendado"] == "upgrade_plano_atual":
            return "media"
        return "baixa"

    df["prioridade_oferta"] = df.apply(prioridade, axis=1)

    # Resumo gerencial
    resumo = (
        df.groupby(["seguro_recomendado", "prioridade_oferta"])
        .agg(
            qtde_clientes=("id_cliente", "count"),
            renda_media=("renda_mensal", "mean"),
            saldo_medio=("saldo_medio", "mean"),
        )
        .reset_index()
    )

    return df, resumo

# GERANDO MENSAGEM PERSONALIZADA COM CO-PILOT

def gerar_prompt_copilot(row: pd.Series) -> str:
    """
    Gera um texto de instrução para o Copilot criar a mensagem de oferta.
    Não chama IA nenhuma aqui: é só o "briefing" que o Copilot vai usar.
    """
    nome = row["nome"]
    idade = row["idade"]
    renda = row["renda_mensal"]
    saldo = row["saldo_medio"]
    faixa_renda = row["faixa_renda"]
    seguro = row["seguro_recomendado"]
    canal = row.get("canal_preferido", "APP")
    prioridade = row["prioridade_oferta"]

    instrucao = (
        f"Você é um gerente de relacionamento do Banco do Ricardo. "
        f"Crie uma mensagem de oferta de seguro personalizada para o cliente {nome}.\n\n"
        f"DADOS DO CLIENTE:\n"
        f"- Idade: {idade}\n"
        f"- Renda mensal: R$ {renda:.2f}\n"
        f"- Saldo médio: R$ {saldo:.2f}\n"
        f"- Faixa de renda: {faixa_renda}\n"
        f"- Seguro recomendado: {seguro}\n"
        f"- Prioridade da oferta: {prioridade}\n"
        f"- Canal preferido de contato: {canal}\n\n"
        f"INSTRUÇÕES PARA A MENSAGEM:\n"
        f"- Escreva em português, com tom profissional e amigável.\n"
        f"- Mencione o Banco do Ricardo pelo nome.\n"
        f"- Destaque os benefícios do tipo de seguro recomendado ({seguro}) de forma simples.\n"
        f"- Sugira que o cliente continue o atendimento pelo canal preferido ({canal}).\n"
        f"- Máximo de 3 a 4 frases."
    )

    return instrucao

def adicionar_prompts_copilot(df: pd.DataFrame) -> pd.DataFrame:
    """Adiciona a coluna 'prompt_copilot' com instruções para o MS Copilot."""
    df = df.copy()
    df["prompt_copilot"] = df.apply(gerar_prompt_copilot, axis=1)
    return df

#iNSTRUCAO DE ENVIO

def definir_instrucao_envio(row: pd.Series) -> str:
    """
    Gera uma descrição ilustrativa de como a mensagem SERIA enviada
    para o canal preferido do cliente, sem de fato enviar nada.
    """
    nome = row["nome"]
    canal = row.get("canal_preferido", "APP")
    prioridade = row["prioridade_oferta"]

    if canal == "APP":
        return (
            f"Enviar notificação no app do Banco do Ricardo para {nome} "
            f"(prioridade {prioridade})."
        )
    elif canal == "WHATSAPP":
        return (
            f"Enviar mensagem via WhatsApp corporativo para {nome} "
            f"(prioridade {prioridade})."
        )
    elif canal == "AGENCIA":
        return (
            f"Incluir {nome} na lista de clientes a serem contatados pelo gerente "
            f"na agência (prioridade {prioridade})."
        )
    elif canal == "EMAIL":
        return (
            f"Enviar e-mail de oferta de seguro para {nome} "
            f"(prioridade {prioridade})."
        )
    else:
        # Caso apareça algum canal diferente
        return (
            f"Agendar contato pelo canal padrão do Banco do Ricardo com {nome} "
            f"(prioridade {prioridade})."
        )


def adicionar_instrucao_envio(df: pd.DataFrame) -> pd.DataFrame:
    """Adiciona a coluna 'instrucao_envio' com a forma ilustrativa de envio."""
    df = df.copy()
    df["instrucao_envio"] = df.apply(definir_instrucao_envio, axis=1)
    return df

# LOAD (l DO ETL) SQLITE + CSV

def load_to_sqlite(df_detalhado: pd.DataFrame, df_resumo: pd.DataFrame):
    """Carrega os dados tratados para um banco SQLite."""
    conn = sqlite3.connect(DB_PATH)
    df_detalhado.to_sql("clientes_banco_do_ricardo", conn, if_exists="replace", index=False)
    df_resumo.to_sql("resumo_seguros_banco_do_ricardo", conn, if_exists="replace", index=False)
    conn.close()


def salvar_csvs(df_detalhado: pd.DataFrame, df_resumo: pd.DataFrame):
    """Salva CSVs de saída com timestamp na pasta processed."""
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    det_path = os.path.join(PROCESSED_DIR, f"clientes_banco_do_ricardo_{ts}.csv")
    res_path = os.path.join(PROCESSED_DIR, f"resumo_seguros_{ts}.csv")

    df_detalhado.to_csv(det_path, index=False)
    df_resumo.to_csv(res_path, index=False)

    print("Arquivos gerados:")
    print(" -", det_path)
    print(" -", res_path)

def salvar_excel(df_detalhado: pd.DataFrame, df_resumo: pd.DataFrame):
    """
    Gera um arquivo Excel com duas abas:
    - 'Clientes'         -> dados detalhados
    - 'Resumo_Seguros'   -> agregados por tipo de seguro
    """
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    xlsx_path = os.path.join(PROCESSED_DIR, f"banco_do_ricardo_ofertas_{ts}.xlsx")

    with pd.ExcelWriter(xlsx_path, engine="xlsxwriter") as writer:
        df_detalhado.to_excel(writer, sheet_name="Clientes", index=False)
        df_resumo.to_excel(writer, sheet_name="Resumo_Seguros", index=False)

        workbook = writer.book
        worksheet_clientes = writer.sheets["Clientes"]

        for i, col in enumerate(df_detalhado.columns):
            max_len = max(12, min(40, df_detalhado[col].astype(str).map(len).max()))
            worksheet_clientes.set_column(i, i, max_len + 2)

    print("Arquivo Excel gerado:")
    print(" -", xlsx_path)    

def gerar_apresentacao(df_detalhado: pd.DataFrame, df_resumo: pd.DataFrame):
    """
    Gera uma apresentação PPTX simples do projeto:
    - Slide de título
    - Slide de visão geral
    - Slide com tabela resumo de seguros
    """
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    pptx_path = os.path.join(PROCESSED_DIR, f"banco_do_ricardo_pipeline_{ts}.pptx")

    prs = Presentation()

    # --- Slide 1: título ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Pipeline de ETL – Banco do Ricardo"
    subtitle.text = (
        "Projeto de recomendação de seguros\n"
        "Geração de prompts via Copilot e orquestração de canais"
    )

    # --- Slide 2: visão geral ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Visão Geral do Pipeline"
    body.text = (
        "1. Extract: leitura de dados do CSV.\n"
        "2. Transform: classificação, regras de seguro e prioridades.\n"
        "3. Prompt Copilot para gerar mensagens personalizadas.\n"
        "4. Canal indicado para envio ilustrativo.\n"
        "5. Load: CSV, SQLite, Excel e apresentação PPTX."
    )

    # --- Slide 3: tabela resumo ---
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Resumo por Tipo de Seguro"

    rows = len(df_resumo) + 1
    cols = 4
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    table.cell(0, 0).text = "Seguro"
    table.cell(0, 1).text = "Prioridade"
    table.cell(0, 2).text = "Qtd. Clientes"
    table.cell(0, 3).text = "Renda Média (R$)"

    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True

    for i, row in enumerate(df_resumo.itertuples(), start=1):
        table.cell(i, 0).text = str(row.seguro_recomendado)
        table.cell(i, 1).text = str(row.prioridade_oferta)
        table.cell(i, 2).text = str(row.qtde_clientes)
        table.cell(i, 3).text = f"{row.renda_media:.2f}"

    prs.save(pptx_path)

    print("Apresentação PPTX gerada:")
    print(" -", pptx_path)    


#PIPELINE

def run_pipeline():
    print("=== E – Extract (Banco do Ricardo) ===")
    df_raw = extract_clientes(RAW_PATH)
    print(f"Clientes lidos: {len(df_raw)}")

    print("=== T – Transform (Regras de Seguros) ===")
    df_enriquecido, resumo = transform(df_raw)

    print("=== T – Prompts para Copilot ===")
    df_enriquecido = adicionar_prompts_copilot(df_enriquecido)

    print("=== T – Instrução ilustrativa de envio por canal ===")
    df_enriquecido = adicionar_instrucao_envio(df_enriquecido)

    print("=== L – Load (SQLite + CSV) ===")
    load_to_sqlite(df_enriquecido, resumo)
    salvar_csvs(df_enriquecido, resumo)

    print("=== EXTRA – Geração de Excel (.xlsx) ===")
    salvar_excel(df_enriquecido, resumo)

    print("=== EXTRA – Geração de apresentação PPTX ===")
    gerar_apresentacao(df_enriquecido, resumo)

    print("Pipeline Banco do Ricardo concluído com sucesso! ✅")


if __name__ == "__main__":
    run_pipeline()

